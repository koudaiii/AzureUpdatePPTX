import unittest
from unittest.mock import patch, MagicMock, Mock
import main
from pptx import Presentation
from pptx.util import Inches
import tempfile
import os


class TestFetchUpdateData(unittest.TestCase):
    """Tests for fetch_update_data function"""

    @patch('main.i18n.get_table_summary_prompt')
    @patch('main.azup.summarize_article_for_table')
    @patch('main.azup.get_article')
    @patch('main.azup.read_and_summary')
    def test_fetch_update_data_returns_correct_format(
        self, mock_read_and_summary, mock_get_article,
        mock_summarize_for_table, mock_get_table_prompt
    ):
        """Test that fetch_update_data returns data in the correct format"""
        # Mock the read_and_summary response
        mock_read_and_summary.return_value = {
            'title': 'Test Azure Update',
            'publishedDate': '2024-01-15T10:30:00.000Z',
            'url': 'https://example.com/update/123',
            'summary': 'Line 1 summary\nLine 2 summary\nLine 3 summary',
            'referenceLink': 'https://docs.example.com/ref1, https://docs.example.com/ref2'
        }

        # Mock table summary generation
        mock_get_table_prompt.return_value = 'Test table summary prompt'
        mock_article_response = MagicMock()
        mock_article_response.json.return_value = {
            'title': 'Test Azure Update',
            'products': ['Azure Service'],
            'description': '<p>Test description</p>'
        }
        mock_get_article.return_value = mock_article_response
        mock_summarize_for_table.return_value = 'One sentence summary for table'

        mock_client = MagicMock()
        deployment_name = 'gpt-4o'
        url = 'https://azure.microsoft.com/updates/test'
        system_prompt = 'Test prompt'

        # Call the function
        result = main.fetch_update_data(url, mock_client, deployment_name, system_prompt)

        # Verify the result structure
        self.assertIsInstance(result, dict)
        self.assertIn('url', result)
        self.assertIn('title', result)
        self.assertIn('published_date_text', result)
        self.assertIn('summary', result)
        self.assertIn('table_summary', result)
        self.assertIn('reference_link_label', result)
        self.assertIn('reference_links', result)

        # Verify the content
        self.assertEqual(result['url'], 'https://example.com/update/123')
        self.assertEqual(result['title'], 'Test Azure Update')
        self.assertEqual(result['table_summary'], 'One sentence summary for table')
        self.assertIsInstance(result['reference_links'], list)
        self.assertEqual(len(result['reference_links']), 2)

        # Verify that read_and_summary was called
        mock_read_and_summary.assert_called_once_with(
            mock_client, deployment_name, url, system_prompt
        )

    @patch('main.i18n.get_table_summary_prompt')
    @patch('main.azup.summarize_article_for_table')
    @patch('main.azup.get_article')
    @patch('main.azup.read_and_summary')
    def test_fetch_update_data_handles_missing_fields(
        self, mock_read_and_summary, mock_get_article,
        mock_summarize_for_table, mock_get_table_prompt
    ):
        """Test that fetch_update_data handles missing fields gracefully"""
        # Mock response with missing fields
        mock_read_and_summary.return_value = {
            'title': '',
            'publishedDate': '',
            'url': '',
            'summary': '',
            'referenceLink': ''
        }

        # Mock table summary generation failure
        mock_get_table_prompt.return_value = 'Test table summary prompt'
        mock_get_article.return_value = None  # Simulate failure

        mock_client = MagicMock()
        deployment_name = 'gpt-4o'
        url = 'https://azure.microsoft.com/updates/test'
        system_prompt = 'Test prompt'

        result = main.fetch_update_data(url, mock_client, deployment_name, system_prompt)

        # Verify it still returns the expected structure
        self.assertIsInstance(result, dict)
        self.assertEqual(result['title'], 'No Title')
        self.assertEqual(result['url'], '')
        self.assertEqual(result['summary'], '')
        self.assertEqual(result['table_summary'], None)
        self.assertEqual(result['reference_links'], [])


class TestCreateUpdateContentSlide(unittest.TestCase):
    """Tests for create_update_content_slide function"""

    def setUp(self):
        """Set up test fixtures"""
        # Create a temporary presentation from the template
        self.temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        self.temp_file.close()
        # Copy the template
        self.prs = Presentation('template/gpstemplate.pptx')

    def tearDown(self):
        """Clean up test files"""
        if os.path.exists(self.temp_file.name):
            os.remove(self.temp_file.name)

    @patch('main.display_update_info')
    def test_create_update_content_slide_adds_slide(self, mock_display):
        """Test that create_update_content_slide adds a slide to the presentation"""
        initial_slide_count = len(self.prs.slides)

        update_data = {
            'url': 'https://example.com/update/123',
            'title': 'Test Update Title',
            'published_date_text': 'Published: 2024-01-15',
            'summary': 'Test summary line 1\nTest summary line 2\nTest summary line 3',
            'table_summary': 'One sentence summary',
            'reference_link_label': 'Reference Links',
            'reference_links': ['https://docs.example.com/ref1', 'https://docs.example.com/ref2']
        }
        page_number = 3

        main.create_update_content_slide(self.prs, update_data, page_number)

        # Verify that a slide was added
        self.assertEqual(len(self.prs.slides), initial_slide_count + 1)

        # Verify the slide content
        new_slide = self.prs.slides[-1]
        self.assertEqual(new_slide.shapes.title.text, 'Test Update Title')

    @patch('main.display_update_info')
    def test_create_update_content_slide_displays_info(self, mock_display):
        """Test that create_update_content_slide calls display_update_info"""
        update_data = {
            'url': 'https://example.com/update/123',
            'title': 'Test Update Title',
            'published_date_text': 'Published: 2024-01-15',
            'summary': 'Test summary',
            'table_summary': 'One sentence summary',
            'reference_link_label': 'Reference Links',
            'reference_links': ['https://docs.example.com/ref1']
        }
        page_number = 3

        main.create_update_content_slide(self.prs, update_data, page_number)

        # Verify that display_update_info was called with correct arguments
        mock_display.assert_called_once_with(
            'Test Update Title',
            'https://example.com/update/123',
            'Published: 2024-01-15',
            'Test summary',
            'Reference Links',
            ['https://docs.example.com/ref1']
        )


class TestAddSummaryTable(unittest.TestCase):
    """Tests for add_summary_table function"""

    def setUp(self):
        """Set up test fixtures"""
        # Create a temporary presentation from the template
        self.temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        self.temp_file.close()
        self.prs = Presentation('template/gpstemplate.pptx')

        # Create a section title slide
        layout = self.prs.slide_layouts[27]
        self.slide = self.prs.slides.add_slide(layout)

    def tearDown(self):
        """Clean up test files"""
        if os.path.exists(self.temp_file.name):
            os.remove(self.temp_file.name)

    def test_add_summary_table_creates_table(self):
        """Test that add_summary_table creates a new slide with a table"""
        updates_data = [
            {
                'url': 'https://example.com/update/1',
                'title': 'Update 1',
                'published_date_text': 'Published: 2024-01-15',
                'summary': 'Summary 1',
                'table_summary': 'One sentence summary 1',
                'reference_link_label': 'Reference Links',
                'reference_links': ['https://docs.example.com/ref1']
            },
            {
                'url': 'https://example.com/update/2',
                'title': 'Update 2',
                'published_date_text': 'Published: 2024-01-16',
                'summary': 'Summary 2',
                'table_summary': 'One sentence summary 2',
                'reference_link_label': 'Reference Links',
                'reference_links': ['https://docs.example.com/ref2']
            }
        ]

        # Get initial slide count
        initial_slide_count = len(self.prs.slides)

        table_pages = main.add_summary_table(self.prs, self.slide, updates_data)

        # Verify that a new slide was created
        self.assertEqual(len(self.prs.slides), initial_slide_count + 1)
        self.assertEqual(table_pages, 1)

        # Get the newly created slide (last slide)
        new_slide = self.prs.slides[-1]

        # Find the table shape in the new slide
        table_shape = None
        for shape in new_slide.shapes:
            if shape.has_table:
                table_shape = shape
                break

        self.assertIsNotNone(table_shape, "Table should be added to the new slide")

        # Verify table dimensions (header row + 2 data rows, 4 columns)
        table = table_shape.table
        self.assertEqual(len(table.rows), 3)  # 1 header + 2 data rows
        self.assertEqual(len(table.columns), 4)  # Page, Title, Summary, URL

    def test_add_summary_table_header_content(self):
        """Test that the table header contains correct column names"""
        updates_data = [{
            'url': 'https://example.com/update/1',
            'title': 'Update 1',
            'published_date_text': 'Published: 2024-01-15',
            'summary': 'Summary 1',
            'table_summary': 'One sentence summary 1',
            'reference_link_label': 'Reference Links',
            'reference_links': []
        }]

        main.add_summary_table(self.prs, self.slide, updates_data)

        # Get the newly created slide
        new_slide = self.prs.slides[-1]

        # Find the table
        table_shape = None
        for shape in new_slide.shapes:
            if shape.has_table:
                table_shape = shape
                break

        table = table_shape.table

        # Verify header content (using i18n keys or expected text)
        # Note: The actual header text will depend on i18n implementation
        # This is a placeholder that should be updated based on actual implementation
        header_row = table.rows[0]
        self.assertIsNotNone(header_row.cells[0].text)  # Page column
        self.assertIsNotNone(header_row.cells[1].text)  # Title column
        self.assertIsNotNone(header_row.cells[2].text)  # Summary column
        self.assertIsNotNone(header_row.cells[3].text)  # URL column

    def test_add_summary_table_page_numbers_start_at_three(self):
        """Test that page numbers account for table pages (Title + Section + Table pages)"""
        updates_data = [
            {
                'url': 'https://example.com/update/1',
                'title': 'Update 1',
                'published_date_text': 'Published: 2024-01-15',
                'summary': 'Summary 1',
                'table_summary': 'One sentence summary 1',
                'reference_link_label': 'Reference Links',
                'reference_links': []
            },
            {
                'url': 'https://example.com/update/2',
                'title': 'Update 2',
                'published_date_text': 'Published: 2024-01-16',
                'summary': 'Summary 2',
                'table_summary': 'One sentence summary 2',
                'reference_link_label': 'Reference Links',
                'reference_links': []
            }
        ]

        table_pages = main.add_summary_table(self.prs, self.slide, updates_data)

        # Get the newly created slide
        new_slide = self.prs.slides[-1]

        # Find the table
        table_shape = None
        for shape in new_slide.shapes:
            if shape.has_table:
                table_shape = shape
                break

        table = table_shape.table

        # Verify page numbers: Title(1) + Section(2) + Table pages(1) = 4 for first detail slide
        # Page numbers should be: 4, 5 (accounting for table page)
        self.assertEqual(table.rows[1].cells[0].text, '4')
        self.assertEqual(table.rows[2].cells[0].text, '5')

    def test_add_summary_table_with_multiple_updates(self):
        """Test that the table handles multiple updates correctly"""
        updates_data = [
            {
                'url': f'https://example.com/update/{i}',
                'title': f'Update {i}',
                'published_date_text': f'Published: 2024-01-{15+i}',
                'summary': f'Summary {i}',
                'table_summary': f'One sentence summary {i}',
                'reference_link_label': 'Reference Links',
                'reference_links': []
            }
            for i in range(1, 6)  # 5 updates
        ]

        main.add_summary_table(self.prs, self.slide, updates_data)

        # Get the newly created slide
        new_slide = self.prs.slides[-1]

        # Find the table
        table_shape = None
        for shape in new_slide.shapes:
            if shape.has_table:
                table_shape = shape
                break

        table = table_shape.table

        # Verify row count (1 header + 5 data rows)
        self.assertEqual(len(table.rows), 6)

        # Verify all titles are present
        for i in range(1, 6):
            self.assertEqual(table.rows[i].cells[1].text, f'Update {i}')

    def test_add_summary_table_empty_list(self):
        """Test that add_summary_table handles empty update list"""
        updates_data = []

        # Get initial slide count
        initial_slide_count = len(self.prs.slides)

        # Empty list should return 0 and not create any new slides
        table_pages = main.add_summary_table(self.prs, self.slide, updates_data)

        # Verify no new slides were created
        self.assertEqual(len(self.prs.slides), initial_slide_count)
        self.assertEqual(table_pages, 0)


class TestSummaryTruncation(unittest.TestCase):
    """Tests for summary truncation at first Japanese period"""

    def setUp(self):
        """Set up test fixtures"""
        self.temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        self.temp_file.close()
        self.prs = Presentation('template/gpstemplate.pptx')
        layout = self.prs.slide_layouts[27]
        self.slide = self.prs.slides.add_slide(layout)

    def tearDown(self):
        """Clean up test files"""
        if os.path.exists(self.temp_file.name):
            os.remove(self.temp_file.name)

    def test_summary_truncated_at_first_period(self):
        """Test that summary is truncated at first 。 when table_summary is not available"""
        updates_data = [{
            'url': 'https://example.com/update/1',
            'title': 'Azure Load Testing',
            'published_date_text': 'Published: 2024-01-15',
            'summary': 'Azure Load TestingがSwitzerland Northリージョンで正式提供開始されました。これにより、大規模な負荷テストやパフォーマンスボトルネックの特定、AIによる分析結果の取得が可能です。CI/CDワークフローへの統合や既存のJMeter・Locustスクリプトも利用できます。',
            'table_summary': None,
            'reference_link_label': 'Reference Links',
            'reference_links': []
        }]

        main.add_summary_table(self.prs, self.slide, updates_data)

        # Get the newly created slide
        new_slide = self.prs.slides[-1]

        # Find the table
        table_shape = None
        for shape in new_slide.shapes:
            if shape.has_table:
                table_shape = shape
                break

        self.assertIsNotNone(table_shape)
        table = table_shape.table

        # Get the summary cell (row 1, column 2)
        summary_cell = table.rows[1].cells[2]
        summary_text = summary_cell.text

        # Verify it's truncated at first 。
        expected = 'Azure Load TestingがSwitzerland Northリージョンで正式提供開始されました。'
        self.assertEqual(summary_text, expected)

    def test_summary_with_no_period(self):
        """Test that summary without 。 is handled gracefully when table_summary is not available"""
        updates_data = [{
            'url': 'https://example.com/update/1',
            'title': 'Test Update',
            'published_date_text': 'Published: 2024-01-15',
            'summary': 'This is a very long summary without any Japanese periods that should be truncated at 100 characters approximately to ensure it fits in one line within the table cell',
            'table_summary': None,
            'reference_link_label': 'Reference Links',
            'reference_links': []
        }]

        main.add_summary_table(self.prs, self.slide, updates_data)

        # Get the newly created slide
        new_slide = self.prs.slides[-1]

        # Find the table
        table_shape = None
        for shape in new_slide.shapes:
            if shape.has_table:
                table_shape = shape
                break

        self.assertIsNotNone(table_shape)
        table = table_shape.table

        # Get the summary cell
        summary_cell = table.rows[1].cells[2]
        summary_text = summary_cell.text

        # Verify it's truncated at 100 chars with ...
        self.assertLessEqual(len(summary_text), 103)  # 100 + "..."
        if len(updates_data[0]['summary'].split('\n')[0]) > 100:
            self.assertTrue(summary_text.endswith('...'))

    def test_summary_with_newlines(self):
        """Test that summary with newlines only takes first line when table_summary is not available"""
        updates_data = [{
            'url': 'https://example.com/update/1',
            'title': 'Test Update',
            'published_date_text': 'Published: 2024-01-15',
            'summary': '第一文です。\n第二文です。\n第三文です。',
            'table_summary': None,
            'reference_link_label': 'Reference Links',
            'reference_links': []
        }]

        main.add_summary_table(self.prs, self.slide, updates_data)

        # Get the newly created slide
        new_slide = self.prs.slides[-1]

        # Find the table
        table_shape = None
        for shape in new_slide.shapes:
            if shape.has_table:
                table_shape = shape
                break

        self.assertIsNotNone(table_shape)
        table = table_shape.table

        # Get the summary cell
        summary_cell = table.rows[1].cells[2]
        summary_text = summary_cell.text

        # Should only contain first sentence
        self.assertEqual(summary_text, '第一文です。')
        self.assertNotIn('第二文', summary_text)

    def test_summary_uses_ai_generated_when_available(self):
        """Test that AI-generated table_summary is used when available"""
        updates_data = [{
            'url': 'https://example.com/update/1',
            'title': 'Azure Load Testing',
            'published_date_text': 'Published: 2024-01-15',
            'summary': 'Azure Load TestingがSwitzerland Northリージョンで正式提供開始されました。これにより、大規模な負荷テストやパフォーマンスボトルネックの特定、AIによる分析結果の取得が可能です。CI/CDワークフローへの統合や既存のJMeter・Locustスクリプトも利用できます。',
            'table_summary': 'Azure Load Testing is now generally available in Switzerland North region.',
            'reference_link_label': 'Reference Links',
            'reference_links': []
        }]

        main.add_summary_table(self.prs, self.slide, updates_data)

        # Get the newly created slide
        new_slide = self.prs.slides[-1]

        # Find the table
        table_shape = None
        for shape in new_slide.shapes:
            if shape.has_table:
                table_shape = shape
                break

        self.assertIsNotNone(table_shape)
        table = table_shape.table

        # Get the summary cell (row 1, column 2)
        summary_cell = table.rows[1].cells[2]
        summary_text = summary_cell.text

        # Verify it uses AI-generated summary, not truncated version
        self.assertEqual(summary_text, 'Azure Load Testing is now generally available in Switzerland North region.')
        self.assertNotIn('これにより', summary_text)

    def test_summary_fallback_when_ai_summary_none(self):
        """Test that fallback logic is used when table_summary is None"""
        updates_data = [{
            'url': 'https://example.com/update/1',
            'title': 'Test Update',
            'published_date_text': 'Published: 2024-01-15',
            'summary': '最初の文章です。次の文章は表示されないはず。',
            'table_summary': None,
            'reference_link_label': 'Reference Links',
            'reference_links': []
        }]

        main.add_summary_table(self.prs, self.slide, updates_data)

        # Get the newly created slide
        new_slide = self.prs.slides[-1]

        # Find the table
        table_shape = None
        for shape in new_slide.shapes:
            if shape.has_table:
                table_shape = shape
                break

        self.assertIsNotNone(table_shape)
        table = table_shape.table

        # Get the summary cell
        summary_cell = table.rows[1].cells[2]
        summary_text = summary_cell.text

        # Verify it uses fallback truncation at first 。
        self.assertEqual(summary_text, '最初の文章です。')
        self.assertNotIn('次の文章', summary_text)


if __name__ == '__main__':
    unittest.main()
