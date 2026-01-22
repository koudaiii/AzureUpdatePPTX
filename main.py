import streamlit as st
import azureupdatehelper as azup
import os
import tempfile
import logging
from pptx import Presentation
from pptx.util import Pt
from datetime import datetime, timedelta
from i18n_helper import i18n, initialize_language_from_query_params
from dotenv import load_dotenv
load_dotenv()

# Initialize language from query parameters before st.set_page_config
initialize_language_from_query_params()

# Set the browser page title
st.set_page_config(
    page_title=i18n.t("page_title"),
    page_icon=":cloud:",
    initial_sidebar_state="auto",
    layout="centered",
    menu_items={
        "Report a bug": "https://github.com/koudaiii/AzureUpdatePPTX/issues",
        "About": f"""### {i18n.t("about_title")}\n{i18n.t("about_content")}""",
    }
)

# Add language selection to sidebar
with st.sidebar:
    i18n.language_selector()

# Set the browser tab title
st.title(i18n.t("main_title"))
# description
st.markdown(i18n.t("description"), unsafe_allow_html=True)

# Today's date (YYYYMMDDHHMMSS) to avoid duplicate file names
save_name = 'AzureUpdates' + datetime.now().strftime('%Y%m%d%H%M%S') + '.pptx'

# Get data from Azure Updates API
entries = azup.get_rss_feed_entries()
st.write(i18n.t(
    "entries_count",
    oldest=azup.oldest_article_date(entries),
    latest=azup.latest_article_date(entries),
    count=len(entries)
))

# Specify how many days back to get updates with streamlit
days = st.slider(i18n.t("slider_label"), 1, 90, 7)


# Set title for Azure Updates slide
def set_slide_title(shape, text, font_size=Pt(24)):
    shape.text = text
    if shape.text_frame and shape.text_frame.paragraphs:
        shape.text_frame.paragraphs[0].font.size = font_size


# Add published_date_text and azure_update_url as one line to Azure Updates slide
def add_hyperlink_text(text_frame, prefix, url, font_size=Pt(18)):
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"{prefix}"
    run.hyperlink.address = url
    run.font.size = font_size


# Add body summary to Azure Updates slide
def add_body_summary(slide, summary):
    body_shape = slide.placeholders[11]
    text_frame = body_shape.text_frame
    text_frame.clear()
    # Add new paragraph if no existing paragraph
    paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    paragraph.text = summary
    paragraph.level = 0


# Add reference links to Azure Updates slide
def add_reference_links(text_frame, label, links):
    # Add header for the reference links
    header = text_frame.add_paragraph()
    header.text = label
    header.level = 2
    # Add each link as a new paragraph with a hyperlink
    for link in links:
        link = link.strip()
        p = text_frame.add_paragraph()
        p.level = 3
        run = p.add_run()
        run.text = link
        run.hyperlink.address = link


# Create title slide
def create_title_slide(prs, title, date_str):
    """
    Creates and configures the title slide using the first layout.

    Args:
        prs: Presentation object.
        title: Title text for the slide.
        date_str: Date string to display in the date placeholder.

    Returns:
        The created slide.
    """
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    # Set the slide title.
    slide.shapes.title.text = title

    # Set the date in the designated placeholder.
    try:
        date_placeholder = slide.placeholders[13]
        date_placeholder.text = date_str
    except IndexError:
        # Log if the expected placeholder index is not found.
        logging.error("Placeholder index 13 for date not found in the title slide.")
    return slide


# Create section title slide
def create_section_title_slide(prs, update_count):
    """
    Creates a section title slide indicating the total number of Azure updates.
    Args:
        prs: The Presentation object.
        update_count: The number of Azure updates.

    Returns:
        A tuple containing the created slide and its first placeholder.
    """
    layout = prs.slide_layouts[27]
    slide = prs.slides.add_slide(layout)

    # Get the title text with proper newline handling
    title_text = i18n.t("section_title", count=update_count)

    # Handle newlines in title by using text frame paragraphs
    title_shape = slide.shapes.title
    title_shape.text = ""  # Clear default text
    text_frame = title_shape.text_frame
    text_frame.clear()  # Clear any existing paragraphs

    # Split text by newlines and create paragraphs
    lines = title_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            # Use the first paragraph
            p = text_frame.paragraphs[0]
        else:
            # Add new paragraphs for additional lines
            p = text_frame.add_paragraph()
        p.text = line

    return slide, slide.placeholders[0]


# Add summary table to a single slide
def add_summary_table_to_slide(slide, updates_data_chunk, start_page_number, font_size=Pt(12)):
    """
    Adds a summary table to a slide.

    Args:
        slide: The slide object to add the table to.
        updates_data_chunk: List of update data dictionaries (subset for this page).
        start_page_number: The starting page number for this chunk.
        font_size: Font size for table content (default: Pt(12)).

    The table contains:
    - Column 1: Page number (starting from 3)
    - Column 2: Title
    - Column 3: Summary (truncated at first 。)
    - Column 4: URL (as hyperlink)
    """
    if not updates_data_chunk:
        return

    # Header row + data rows
    rows = 1 + len(updates_data_chunk)
    cols = 4  # Page, Title, Summary, URL

    # Position the table: 1cm from top (1cm ≈ 28.35pt)
    # More compact layout with smaller font
    left = Pt(40)
    top = Pt(28)
    width = Pt(880)
    height = Pt(540)  # Increased height for more rows

    # Create the table
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Pt(50)   # Page number - narrow
    table.columns[1].width = Pt(220)  # Title - medium
    table.columns[2].width = Pt(380)  # Summary - wide
    table.columns[3].width = Pt(230)  # URL - medium

    # Set row heights to make table more compact
    for row in table.rows:
        row.height = Pt(30)  # Compact row height

    # Set header row with unified font size
    header_cells = table.rows[0].cells
    header_cells[0].text = i18n.t("table_header_page")
    header_cells[1].text = i18n.t("table_header_title")
    header_cells[2].text = i18n.t("table_header_summary")
    header_cells[3].text = i18n.t("table_header_url")

    # Style header row
    for cell in header_cells:
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = font_size
        # Reduce cell margins for more compact layout
        cell.text_frame.margin_top = Pt(2)
        cell.text_frame.margin_bottom = Pt(2)
        cell.text_frame.margin_left = Pt(5)
        cell.text_frame.margin_right = Pt(5)

    # Fill data rows
    for idx, update_data in enumerate(updates_data_chunk):
        row_idx = idx + 1  # Skip header row
        page_number = start_page_number + idx

        cells = table.rows[row_idx].cells

        # Page number
        cells[0].text = str(page_number)
        cells[0].text_frame.paragraphs[0].font.size = font_size
        cells[0].text_frame.margin_top = Pt(2)
        cells[0].text_frame.margin_bottom = Pt(2)
        cells[0].text_frame.margin_left = Pt(5)
        cells[0].text_frame.margin_right = Pt(5)

        # Title
        cells[1].text = update_data['title']
        cells[1].text_frame.paragraphs[0].font.size = font_size
        cells[1].text_frame.margin_top = Pt(2)
        cells[1].text_frame.margin_bottom = Pt(2)
        cells[1].text_frame.margin_left = Pt(5)
        cells[1].text_frame.margin_right = Pt(5)

        # Summary - use AI-generated one-sentence summary if available, otherwise fallback
        if update_data.get('table_summary'):
            # Use AI-generated one-sentence summary
            summary_text = update_data['table_summary']
            logging.debug(f"Using AI-generated table summary for: {update_data['title']}")
        else:
            # Fallback: truncate at first 。 (Japanese period)
            summary = update_data['summary']
            # Take only the first line
            first_line = summary.split('\n')[0] if summary else ""
            # Find the first 。 and truncate there (one sentence only)
            if '。' in first_line:
                first_period_idx = first_line.find('。')
                summary_text = first_line[:first_period_idx + 1]  # Include the 。
            else:
                # If no 。, limit to 100 characters
                summary_text = first_line[:100] + "..." if len(first_line) > 100 else first_line
            logging.debug(f"Using fallback truncation for table summary: {update_data['title']}")
        cells[2].text = summary_text
        cells[2].text_frame.paragraphs[0].font.size = font_size
        cells[2].text_frame.margin_top = Pt(2)
        cells[2].text_frame.margin_bottom = Pt(2)
        cells[2].text_frame.margin_left = Pt(5)
        cells[2].text_frame.margin_right = Pt(5)

        # URL with hyperlink
        url = update_data['url']
        url_cell = cells[3]
        url_cell.text_frame.clear()
        url_cell.text_frame.margin_top = Pt(2)
        url_cell.text_frame.margin_bottom = Pt(2)
        url_cell.text_frame.margin_left = Pt(5)
        url_cell.text_frame.margin_right = Pt(5)
        p = url_cell.text_frame.paragraphs[0]
        run = p.add_run()
        # Display shortened URL text but link to full URL
        display_text = url if len(url) <= 45 else url[:42] + "..."
        run.text = display_text
        run.hyperlink.address = url
        run.font.size = font_size


# Add summary tables to presentation (with pagination support)
def add_summary_table(prs, section_slide, updates_data, max_rows_per_page=8):
    """
    Adds summary table(s) to the presentation using layout 28 (blank), splitting into multiple slides if needed.
    The section_slide parameter is kept for compatibility but not used (all tables use layout 28).

    Args:
        prs: The Presentation object.
        section_slide: The section title slide (kept for compatibility, not used for tables).
        updates_data: List of all update data dictionaries.
        max_rows_per_page: Maximum number of data rows per page (default: 8).

    Returns:
        Number of table slides created.
    """
    if not updates_data:
        return 0

    total_updates = len(updates_data)
    pages_needed = (total_updates + max_rows_per_page - 1) // max_rows_per_page

    for page_idx in range(pages_needed):
        start_idx = page_idx * max_rows_per_page
        end_idx = min(start_idx + max_rows_per_page, total_updates)
        chunk = updates_data[start_idx:end_idx]
        # Page numbers for detail slides: Title(1) + Section(2) + Table pages + offset
        start_page_number = 3 + pages_needed + start_idx

        # All table pages use layout 28 (blank layout for table only)
        layout = prs.slide_layouts[28]
        new_slide = prs.slides.add_slide(layout)
        logging.info(f"Creating table slide {page_idx + 1} of {pages_needed} using layout 28 (blank)")

        # No title needed - table only
        add_summary_table_to_slide(new_slide, chunk, start_page_number)

    return pages_needed


# Display Azure Updates information
def display_update_info(title, url, published_date, summary, ref_label, ref_links):
    st.write('')
    st.markdown(f"「{title}」", unsafe_allow_html=True)
    st.markdown(
        f"<small><a href='{url}' target='_blank'>{published_date}</a></small>",
        unsafe_allow_html=True
    )
    st.markdown(f"<small>{summary}</small>", unsafe_allow_html=True)
    st.markdown(f"<small>{ref_label}</small>", unsafe_allow_html=True)
    for link in ref_links:
        st.markdown(
            f"<p style='line-height:130%; margin-bottom:10px; padding:0;'>"
            f"<a href='{link}' target='_blank'>{link}</a></p>",
            unsafe_allow_html=True
        )
    st.write('')


# Add to Azure Updates slide
def create_update_slide(prs, title, published_date, url, summary, ref_label, ref_links):
    """Creates a new slide for an Azure Updates and configures its elements."""
    layout = prs.slide_layouts[10]
    slide = prs.slides.add_slide(layout)

    # Set slide title
    set_slide_title(slide.shapes.title, title)

    # Add published date with hyperlink (using placeholder index 10)
    try:
        ph_date = slide.placeholders[10].text_frame
        add_hyperlink_text(ph_date, published_date, url)
    except IndexError:
        logging.error("Placeholder index 10 not found in update slide.")

    # Add summary/body content
    add_body_summary(slide, summary)

    # Add reference links (using placeholder index 11)
    try:
        ph_refs = slide.placeholders[11].text_frame
        add_reference_links(ph_refs, ref_label, ref_links)
    except IndexError:
        logging.error("Placeholder index 11 not found in update slide.")

    return slide


# Generate Azure Updates content
def extract_update_data(result):
    title = result.get("title", "No Title") or "No Title"  # Handle empty string
    published_date_raw = result.get("publishedDate", "")
    published_date_str = published_date_raw.split(".")[0] if published_date_raw else ""
    try:
        dt = datetime.strptime(published_date_str, '%Y-%m-%dT%H:%M:%S')
        published_date_text = i18n.t("published_date", date=i18n.format_date(dt))
    except ValueError:
        published_date_text = i18n.t("published_date", date="Unknown")
    url = result.get("url", "")
    summary = result.get("summary", "")
    ref_label = i18n.t("reference_links")
    ref_links = [link.strip() for link in result.get("referenceLink", "").split(",") if link.strip()]
    return title, published_date_text, url, summary, ref_label, ref_links


# Create Azure Updates
def process_update(url, client, deployment_name, prs, system_prompt):
    # Process and log Azure Updates information
    logging.info("***** Begin of Record *****")
    result = azup.read_and_summary(client, deployment_name, url, system_prompt)
    logging.debug("Result: %s", result)
    for key, value in result.items():
        logging.info("%s : %s", key, value)
    logging.info("***** End of Record *****")

    # Extract update data from the result
    (
        azure_update_title,
        published_date_text,
        azure_update_url,
        azure_update_summary,
        reference_link_label,
        reference_links,
    ) = extract_update_data(result)

    # Display update information via Streamlit
    display_update_info(azure_update_title, azure_update_url, published_date_text,
                        azure_update_summary, reference_link_label, reference_links)

    # Create and add the update slide to the presentation
    create_update_slide(prs, azure_update_title, published_date_text, azure_update_url,
                        azure_update_summary, reference_link_label, reference_links)


# Title slide title
def generate_slide_info(start_date, end_date) -> tuple[str, str]:
    slide_title = i18n.t(
        "slide_title",
        start=start_date.strftime('%Y/%m/%d'),
        end=end_date.strftime('%Y/%m/%d')
    )
    return slide_title


# Display list of Azure Updates URLs
def display_update_urls(urls):
    update_count = len(urls)
    st.write(i18n.t("update_count", count=update_count))
    st.write(urls)


# Specified date
def start_date(days):
    return datetime.now().astimezone() - timedelta(days)


# End time is current time
def end_date():
    return datetime.now().astimezone()


# Press button to get data from Azure Updates API and generate PPTX
if st.button(i18n.t("button_text")):
    # Display error and exit if environment variables are missing
    if not azup.environment_check():
        st.error(i18n.t("env_error"))
        st.stop()

    st.write(i18n.t(
        "date_range",
        start=start_date(days).strftime("%Y-%m-%d"),
        end=end_date().strftime("%Y-%m-%d")
    ))
    urls = azup.target_update_urls(entries, start_date(days))
    display_update_urls(urls)

    # PPTX generation process
    st.write(i18n.t("generating"))

    # Generate slide title and date string
    slide_title = generate_slide_info(start_date(days), end_date())

    # Create a temporary PPTX file using a template
    pptx_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs = Presentation("template/gpstemplate.pptx")

    # First slide (title slide)
    slide = create_title_slide(prs, slide_title, end_date().strftime('%Y%m%d%H%M%S'))
    # Second slide (section title slide)
    slide, date_ph = create_section_title_slide(prs, len(urls))

    # Get Azure OpenAI client
    client, deployment_name = azup.azure_openai_client(os.getenv("API_KEY"), os.getenv("API_ENDPOINT"))

    # system prompt for Azure OpenAI
    system_prompt = i18n.get_system_prompt()

    # Get Azure Updates information and create slides
    for url in urls:
        process_update(url, client, deployment_name, prs, system_prompt)

    # Save PPTX
    prs.save(pptx_file.name)
    st.write(i18n.t("done"))

    try:
        with open(pptx_file.name, "rb") as f:
            st.download_button(i18n.t("download_button"), f.read(), file_name=save_name)
    finally:
        pptx_file.close()
        # Delete temporary file
        os.remove(pptx_file.name)
